from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

from .frames import Frame

EMU_PER_PX = 9525  # 96 dpi


def build_pdf(image_paths: list[Path], out_path: Path) -> None:
    images = [Image.open(p).convert("RGB") for p in image_paths]
    if not images:
        raise ValueError("Нет картинок для сборки PDF.")
    images[0].save(out_path, save_all=True, append_images=images[1:])


def build_pptx(frames: list[Frame], image_paths: list[Path], out_path: Path) -> None:
    if not frames or not image_paths:
        raise ValueError("Нет слайдов для сборки PPTX.")

    slide_w = max((f.width for f in frames if f.width), default=1280)
    slide_h = max((f.height for f in frames if f.height), default=720)

    prs = Presentation()
    prs.slide_width = Emu(int(slide_w * EMU_PER_PX))
    prs.slide_height = Emu(int(slide_h * EMU_PER_PX))
    blank_layout = prs.slide_layouts[6]

    for frame, image_path in zip(frames, image_paths):
        slide = prs.slides.add_slide(blank_layout)
        left, top, width, height = _fit(frame.width, frame.height, slide_w, slide_h)
        slide.shapes.add_picture(
            str(image_path),
            Emu(int(left * EMU_PER_PX)),
            Emu(int(top * EMU_PER_PX)),
            width=Emu(int(width * EMU_PER_PX)),
            height=Emu(int(height * EMU_PER_PX)),
        )

    prs.save(out_path)


def _fit(w: float, h: float, box_w: float, box_h: float) -> tuple[float, float, float, float]:
    """Вписывает w×h в box_w×box_h с сохранением пропорций, по центру."""
    if not w or not h:
        return 0.0, 0.0, box_w, box_h
    scale = min(box_w / w, box_h / h)
    fw, fh = w * scale, h * scale
    return (box_w - fw) / 2, (box_h - fh) / 2, fw, fh
